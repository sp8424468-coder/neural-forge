from flask import Flask, render_template, request, redirect, session
from PyPDF2 import PdfReader
from groq import Groq
from pptx import Presentation

from dotenv import load_dotenv
load_dotenv()

from flask_mail import Mail, Message
import random
import os
import sqlite3


app = Flask(__name__)

app.secret_key = os.getenv("SECRET_KEY")

client = Groq(
    api_key=os.getenv("GROQ_API_KEY")
)

app.config['MAIL_SERVER'] = 'smtp.gmail.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USERNAME'] = os.getenv("MAIL_USERNAME")
app.config['MAIL_PASSWORD'] = os.getenv("MAIL_PASSWORD")

mail = Mail(app)
@app.route("/quiz")
def quiz():
    if "user" not in session:
        return redirect("/login")

    return render_template("quiz.html")

@app.route("/send-otp")
def send_otp():

    user = session.get("temp_user")

    if not user:
        return redirect("/register")

    email = user["email"]

    otp = str(random.randint(100000,999999))
    session["otp"] = otp

    msg = Message("Your OTP - MindTrack AI",
                  sender=app.config['MAIL_USERNAME'],
                  recipients=[email])

    msg.body = f"Your OTP is: {otp}"

    mail.send(msg)

    return redirect("/verify-otp")
@app.route("/verify-otp", methods=["GET","POST"])
def verify_otp():

    if request.method == "POST":

        user_otp = request.form["otp"]

        if user_otp == session.get("otp"):

            user = session.get("temp_user")

            if not user:
                return redirect("/register")

            db = get_db()
            cur = db.cursor()

            # ✅ SAVE USER AFTER OTP VERIFIED
            cur.execute(
                "INSERT INTO users(name,email,password) VALUES(?,?,?)",
                (user["name"], user["email"], user["password"])
            )

            db.commit()

            # 🔥 clear temp data
            session.pop("temp_user", None)
            session.pop("otp", None)

            return redirect("/login")

        else:
            return "❌ Invalid OTP"

    return render_template("verify_otp.html")


def get_db():
    return sqlite3.connect("database.db")

# ---------------- HOME ----------------
@app.route("/")
def home():
    return render_template("index.html")


# ---------------- REGISTER ----------------
@app.route("/register", methods=["GET","POST"])
def register():

    if request.method == "POST":

        name = request.form["name"]
        email = request.form["email"]
        password = request.form["password"]

        db = get_db()
        cur = db.cursor()

        # 🔥 CHECK IF EMAIL ALREADY EXISTS (ADD HERE)
        cur.execute("SELECT * FROM users WHERE email=?", (email,))
        existing = cur.fetchone()

        if existing:
            return render_template("register.html", error="Email already registered ❌")

        # 🔥 store data temporarily (ONLY IF NEW USER)
        session["temp_user"] = {
            "name": name,
            "email": email,
            "password": password
        }

        # 🔥 redirect to OTP
        return redirect("/send-otp")

    return render_template("register.html")

# ---------------- LOGIN ----------------

@app.route("/login", methods=["GET","POST"])
def login():

    # 🔥 ALWAYS CLEAR OLD SESSION
    session.clear()

    if request.method == "POST":

        email = request.form["email"].strip()
        password = request.form["password"].strip()

        db = get_db()
        cur = db.cursor()

        user = cur.execute(
            "SELECT name, email, password FROM users WHERE email=?",
            (email,)
        ).fetchone()

        # 🔥 STRICT CHECK
        if user and user[2] == password:

            session["user"] = user[1]   # email
            session["name"] = user[0]   # name

            return redirect("/dashboard")

        else:
            return render_template("login.html", error="❌ Invalid Email or Password")

    return render_template("login.html")


# -----------------------------
# Build HYBRID prompt
# -----------------------------
@app.route("/chat-api", methods=["POST"])
def chat_api():

    message = request.json["message"]
    user = session["user"]

    db = get_db()
    cur = db.cursor()

    # -----------------------------
    # Get notes (RAG)
    # -----------------------------
    rows = cur.execute(
        "SELECT content FROM notes WHERE user=?",
        (user,)
    ).fetchall()

    notes_text = " ".join([r[0] for r in rows])
    context = notes_text[:2000]

    # -----------------------------
    # Get current topic
    # -----------------------------
    topic = session.get("current_topic", "")

    # -----------------------------
    # 🔥 FINAL PROMPT (STRICT FORMAT)
    # -----------------------------
    prompt = f"""
You are a smart study tutor.

Topic: {topic}

Student Notes:
{context}

Question:
{message}

Answer STRICTLY in this format:

Definition:
<one simple line>

Explanation:
<2-3 simple lines>

Example:
<one real-world example, no code>

Rules:
- Do NOT use symbols (*, #, |, ``` )
- Do NOT give code
- Use plain simple English
- Keep each section in new line
"""

    try:
        response = client.chat.completions.create(
            model="openai/gpt-oss-120b",
            messages=[{"role":"user","content":prompt}],
            max_tokens=300
        )

        reply = response.choices[0].message.content

        # -----------------------------
        # 🔥 CLEAN RESPONSE
        # -----------------------------
        reply = (
            reply
            .replace("*", "")
            .replace("#", "")
            .replace("|", "")
            .replace("```", "")
        )

    except Exception as e:
        print(e)
        reply = "⚠️ AI error"

    return {"reply": reply}

@app.route("/generate-quiz", methods=["POST"])
def generate_quiz():

    if "user" not in session:
        return {"quiz": "[]"}

    data = request.json
    topic = data.get("topic")
    difficulty = data.get("difficulty")
    note_id = data.get("note_id")  # 🔥 IMPORTANT (selected file)

    user = session["user"]

    db = get_db()
    cur = db.cursor()

    # -----------------------------
    # 🔥 GET CONTENT FROM SELECTED FILE
    # -----------------------------
    if note_id:
        row = cur.execute(
            "SELECT content FROM notes WHERE id=? AND user=?",
            (note_id, user)
        ).fetchone()

        content = (row[0] if row else "")[:2000]

    else:
        # fallback → use all notes
        rows = cur.execute(
            "SELECT content FROM notes WHERE user=?",
            (user,)
        ).fetchall()

        content = " ".join([r[0] for r in rows])[:2000]

    # -----------------------------
    # ❗ SAFETY (empty content)
    # -----------------------------
    if not content.strip():
        content = "Basic concepts introduction overview definitions examples"

    # -----------------------------
    # 🎯 PROMPTS (CONTENT-BASED)
    # -----------------------------
    if difficulty == "easy":
        prompt = f"""
Create 3 MCQ questions ONLY from the given study content.

Topic: {topic}

Study Content:
{content}

Rules:
- Questions must come ONLY from this content
- Do NOT use outside knowledge
- Keep simple
- 4 options each

Return ONLY JSON:
[
  {{
    "question": "question",
    "options": ["A","B","C","D"],
    "answer": "correct option"
  }}
]
"""

    elif difficulty == "medium":
        prompt = f"""
Create 3 short answer questions ONLY from the given study content.

Topic: {topic}

Study Content:
{content}

Rules:
- Based ONLY on content
- 2-3 line answers

Return ONLY JSON:
[
  {{
    "question": "question",
    "answer": "short answer"
  }}
]
"""

    else:
        prompt = f"""
Create 3 long answer questions ONLY from the given study content.

Topic: {topic}

Study Content:
{content}

Rules:
- Based ONLY on content
- Detailed explanation answers

Return ONLY JSON:
[
  {{
    "question": "question",
    "answer": "detailed answer"
  }}
]
"""

    # -----------------------------
    # 🔥 AI CALL
    # -----------------------------
    try:
        response = client.chat.completions.create(
            model="openai/gpt-oss-120b",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=600
        )

        quiz = response.choices[0].message.content

        return {"quiz": quiz}

    except Exception as e:
        print("AI ERROR:", e)

        # -----------------------------
        # 🔥 FALLBACK (SAFE)
        # -----------------------------
        return {
            "quiz": f"""
[
  {{
    "question": "Basic question about {topic}?",
    "options": ["Concept","Definition","Example","All"],
    "answer": "Concept"
  }},
  {{
    "question": "What is related to {topic}?",
    "options": ["Idea","Method","Process","All"],
    "answer": "All"
  }}
]
"""
        }

@app.route("/submit-quiz", methods=["POST"])
def submit_quiz():

    topic = request.form["topic"]
    score = request.form["score"]

    db = get_db()
    cur = db.cursor()

    cur.execute(
        "CREATE TABLE IF NOT EXISTS analytics(user TEXT, topic TEXT, score INTEGER)"
    )

    cur.execute(
        "INSERT INTO analytics(user,topic,score) VALUES(?,?,?)",
        (session["user"], topic, score)
    )

    db.commit()

    return redirect("/study-plan")



# ---------------- DASHBOARD ----------------
@app.route("/dashboard")
def dashboard():

    if "user" not in session:
        return redirect("/login")

    db = get_db()
    cur = db.cursor()

    user_email = session["user"]

    # 🔥 GET USER NAME
    user = cur.execute(
        "SELECT name FROM users WHERE email=?",
        (user_email,)
    ).fetchone()

    if not user:
        session.clear()
        return redirect("/login")

    name = user[0]

    # ✅ TOTAL TOPICS
    total = cur.execute(
        "SELECT COUNT(*) FROM study_plan WHERE user=?",
        (user_email,)
    ).fetchone()[0]

    # ✅ COMPLETED TOPICS
    completed = cur.execute(
        "SELECT COUNT(*) FROM study_plan WHERE user=? AND status='done'",
        (user_email,)
    ).fetchone()[0]

    # ✅ NOTES COUNT
    notes = cur.execute(
        "SELECT COUNT(*) FROM notes WHERE user=?",
        (user_email,)
    ).fetchone()[0]

    # ✅ QUIZ ACCURACY
    score = cur.execute(
        "SELECT AVG(score) FROM analytics WHERE user=?",
        (user_email,)
    ).fetchone()[0]

    accuracy = int(score) if score else 0

    return render_template(
        "dashboard.html",
        name=name,
        total=total,
        completed=completed,
        notes=notes,
        accuracy=accuracy
    )



@app.route("/notes")
def notes():

    if "user" not in session:
        return redirect("/login")

    db = get_db()
    cur = db.cursor()

    user = session["user"]

    try:
        rows = cur.execute(
            "SELECT id,subject,content FROM notes WHERE user=?",
            (user,)
        ).fetchall()
    except:
        rows = []

    return render_template("notes.html", notes=rows)


@app.route("/save-notes", methods=["POST"])
def save_notes():

    if "user" not in session:
        return redirect("/login")

    subject = request.form.get("subject")
    user = session["user"]

    file = request.files.get("file")
    content = ""

    try:
        # ---------------- PDF ----------------
        if file and file.filename.lower().endswith(".pdf"):

            reader = PdfReader(file)

            for page in reader.pages:
                content += page.extract_text() or ""

        # ---------------- PPT ----------------
        elif file and file.filename.lower().endswith(".pptx"):

            prs = Presentation(file)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + " "

        # ---------------- TEXT INPUT ----------------
        else:
            content = request.form.get("content") or ""

        # ---------------- SAFETY ----------------
        if not content.strip():
            content = "No readable content found"

        # ---------------- SAVE ----------------
        db = get_db()
        cur = db.cursor()

        cur.execute(
            "INSERT INTO notes(user,subject,content) VALUES(?,?,?)",
            (user, subject, content)
        )

        db.commit()

    except Exception as e:
        print("ERROR:", e)
        return "❌ Error processing file"

    return redirect("/notes")
# ---------------- CHAT PAGE ----------------
@app.route("/chat")
def chat():
    if "user" not in session:
        return redirect("/login")

    return render_template("chat.html")

@app.route("/delete-note/<int:id>")
def delete_note(id):

    if "user" not in session:
        return redirect("/login")

    user = session["user"]

    db = get_db()
    cur = db.cursor()

    cur.execute(
        "DELETE FROM notes WHERE id=? AND user=?",
        (id, user)
    )

    db.commit()

    return redirect("/notes")
@app.route("/learn/<topic>")
def learn(topic):

    session["current_topic"] = topic

    return redirect("/chat")

@app.route("/generate-plan", methods=["POST"])
def generate_plan():

    if "user" not in session:
        return redirect("/login")

    user = session["user"]
    note_id = request.form.get("note_id")

    db = get_db()
    cur = db.cursor()

    # ---------------- GET NOTE ----------------
    row = cur.execute(
        "SELECT content FROM notes WHERE id=? AND user=?",
        (note_id, user)
    ).fetchone()

    if not row:
        return "❌ Note not found"

    content = (row[0] or "").strip()

    # ---------------- HANDLE EMPTY CONTENT ----------------
    if not content or content == "No readable content found":
        topics = [
            "Main Concepts",
            "Key Ideas",
            "Important Topics"
        ]

    else:
        content = content[:2000]

        # 🔥 IMPROVED PROMPT (VERY IMPORTANT)
        prompt = f"""
Extract ONLY meaningful study topics from this content.

IMPORTANT:
- Topics must be SPECIFIC to the subject
- Do NOT include generic words like:
  Overview, Introduction, Basics, Summary
- Topics should be real subject concepts

Example:
Wrong ❌: Overview, Basics
Correct ✅: Data Analysis, Multimedia Systems, Machine Learning

Rules:
- Max 8 topics
- One topic per line
- No explanation

Content:
{content}
"""

        try:
            response = client.chat.completions.create(
                model="openai/gpt-oss-120b",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=150
            )

            raw = response.choices[0].message.content

            # ---------------- CLEAN + FILTER ----------------
            topics = [
                t.strip()
                for t in raw.split("\n")
                if t.strip() and len(t.strip()) < 50
            ]

            # 🔥 REMOVE BAD GENERIC TOPICS
            bad_words = ["overview", "introduction", "basics", "summary"]

            topics = [
                t for t in topics
                if t.lower() not in bad_words
            ]

            # 🔥 FINAL SAFETY
            if not topics:
                topics = [
                    "Core Concepts",
                    "Key Topics",
                    "Applications",
                    "Examples"
                ]

        except Exception as e:
            print("AI ERROR:", e)

            topics = [
                "Core Concepts",
                "Applications",
                "Examples"
            ]

    # ---------------- SAVE TOPICS (NO DUPLICATES) ----------------
    for t in topics:
        exists = cur.execute(
            "SELECT 1 FROM study_plan WHERE user=? AND topic=?",
            (user, t)
        ).fetchone()

        if not exists:
            cur.execute(
                "INSERT INTO study_plan(user, topic, status) VALUES(?,?,?)",
                (user, t, "pending")
            )

    db.commit()

    # ---------------- FETCH FINAL DATA ----------------
    rows = cur.execute(
        "SELECT topic, status FROM study_plan WHERE user=?",
        (user,)
    ).fetchall()

    return render_template("study_plan.html", topics=rows)
# ---------------- AI EVALUATE ANSWER ----------------
@app.route("/evaluate-answer", methods=["POST"])
def evaluate_answer():

    data = request.json

    question = data["question"]
    correct_answer = data["answer"]
    student_answer = data["student_answer"]

    prompt = f"""
You are an exam evaluator.

Question:
{question}

Correct Answer:
{correct_answer}

Student Answer:
{student_answer}

Evaluate like a strict teacher.

Rules:
- If answer is correct → Marks: 1
- If wrong → Marks: 0
- Give short feedback (1 line)
- Be simple and clear

Return EXACT format:

Result: correct OR wrong
Marks: 0 or 1
Feedback: <one line>
"""

    try:
        response = client.chat.completions.create(
            model="openai/gpt-oss-120b",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=150
        )

        reply = response.choices[0].message.content

        # 🔥 CHECK RESULT
        if "Marks: 1" in reply:
            result = "correct"

            # ✅ UPDATE STUDY PLAN (VERY IMPORTANT)
            db = get_db()
            cur = db.cursor()

            cur.execute("""
            UPDATE study_plan 
            SET status='done' 
            WHERE user=? AND topic=?
            """, (session["user"], session.get("current_topic","")))

            db.commit()

        else:
            result = "wrong"

        return {
            "result": result,
            "full": reply
        }

    except Exception as e:
        print(e)
        return {
            "result": "wrong",
            "full": "⚠️ AI evaluation error"
        }
@app.route("/set-topic")
def set_topic():

    topic = request.args.get("topic")
    session["current_topic"] = topic

    return redirect("/chat")


@app.route("/get-topic")
def get_topic():
    return {"topic": session.get("current_topic", "")}


@app.route("/study-plan")
def study_plan_page():

    if "user" not in session:
        return redirect("/login")

    db = get_db()
    cur = db.cursor()

    rows = cur.execute(
        "SELECT topic, status FROM study_plan WHERE user=?",
        (session["user"],)
    ).fetchall()

    return render_template("study_plan.html", topics=rows)

# ---------------- LOGOUT ----------------
@app.route("/logout")
def logout():
    session.clear()
    return redirect("/")

if __name__ == "__main__":
    app.run(debug=True)